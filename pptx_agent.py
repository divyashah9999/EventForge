import os
import re
import sys

# Force UTF-8 output on Windows terminals that default to cp1252
if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a hex color string like '#1A1A2E' or '1A1A2E' to RGBColor."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return RGBColor(0x1A, 0x1A, 0x2E)  # safe dark default
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def generate_pitch_deck(json_data: dict) -> str:
    """
    Creates a PowerPoint pitch deck from json_data['pitch_deck'].
    Applies color_palette hex codes to slide titles.
    Saves as {title}_PitchDeck.pptx and returns the file path.
    """
    pitch_deck = json_data.get("pitch_deck", [])
    if not pitch_deck:
        raise ValueError("json_data['pitch_deck'] is empty or missing.")

    # Pull color palette — expect a list of hex strings
    palette = json_data.get("color_palette", [])
    primary_color   = _hex_to_rgb(palette[0]) if len(palette) > 0 else RGBColor(0x1A, 0x1A, 0x2E)
    accent_color    = _hex_to_rgb(palette[1]) if len(palette) > 1 else RGBColor(0x7C, 0x3A, 0xED)
    text_color      = _hex_to_rgb(palette[2]) if len(palette) > 2 else RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[1]  # "Title and Content" — per SYSTEM_PROMPT spec

    for slide_data in pitch_deck:
        slide = prs.slides.add_slide(slide_layout)

        # --- Title ---
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "")
        title_tf = title_shape.text_frame
        for para in title_tf.paragraphs:
            for run in para.runs:
                run.font.size  = Pt(32)
                run.font.bold  = True
                run.font.color.rgb = primary_color

        # --- Content body ---
        body_shape = slide.placeholders[1]
        body_tf = body_shape.text_frame
        body_tf.word_wrap = True
        body_tf.text = slide_data.get("content", "")
        for para in body_tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # --- Visual guidance note (small italic at bottom) ---
        visual_guidance = slide_data.get("visual_guidance", "")
        if visual_guidance:
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5)
            )
            tf = txBox.text_frame
            tf.text = f"🎨 Visual: {visual_guidance}"
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.runs[0]
            run.font.size   = Pt(9)
            run.font.italic = True
            run.font.color.rgb = accent_color

    # Safe filename — strip any characters that aren't alphanumeric, spaces, or hyphens
    safe_title = re.sub(r"[^\w\s-]", "", json_data.get("title", "Event")).strip()
    filename = f"{safe_title.replace(' ', '_')}_PitchDeck.pptx"
    filepath = os.path.join(os.getcwd(), filename)

    prs.save(filepath)
    print(f"✅ Pitch deck saved: {filepath}")
    return filepath
